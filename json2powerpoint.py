import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import random
import ast

questions = ast.literal_eval(open('questions_output.txt', 'r').read())

print('Shuffling questions...')
l = list(questions.items())
random.shuffle(l)
questions = dict(l)

print('Generating quiz...')
pr1 = Presentation()


for x, i in enumerate(questions):
    question_slide_register = pr1.slide_layouts[0]
    question_slide = pr1.slides.add_slide(question_slide_register)

    question_title = question_slide.shapes[0]
    question_subtitle = question_slide.shapes[1]
    question_title.text = f'Frage {x+1}:\n{i}'
    question_title.text_frame.paragraphs[1].font.size = Pt(24)
    question_subtitle.text = '\n\nBitte antworten...'
    question_subtitle.text_frame.paragraphs[2].font.color.rgb = RGBColor(240, 202, 12)


    answer_slide_register = pr1.slide_layouts[0]
    answer_slide = pr1.slides.add_slide(answer_slide_register)

    answer_title = answer_slide.shapes[0]
    answer_subtitle = answer_slide.shapes[1]
    answer_title.text = f'Frage {x+1}:\n{i}'
    answer_title.text_frame.paragraphs[1].font.size = Pt(24)
    answer_subtitle.text = '\n\n' + questions[i]
    answer_subtitle.text_frame.paragraphs[2].font.color.rgb = RGBColor(9, 214, 15)

    print(f'Question {x} added.')

print('PPTX generated. Saving...')
pr1.save('Der dümmste fliegt!.pptx')